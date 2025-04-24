
import translation_agent as ta
import argparse
import time
import os
from docx import Document
from pptx import Presentation
from openpyxl import Workbook, load_workbook
from bs4 import BeautifulSoup
from ebooklib import epub
import re
from md2pdf.core import md2pdf
import markdown
from docx import Document
from bs4 import BeautifulSoup
from PIL import Image, ImageDraw, ImageFont
from surya.recognition import RecognitionPredictor
from surya.detection import DetectionPredictor


def translate_image(image_path, source_lang=None, target_lang='Russian', do_reflection=False, country='Russia'):

    # Шрифт для отрисовки текста
    font_path = "arial.ttf"
    font_size = 16
    font = ImageFont.truetype(font_path, font_size)
    white_color = (255, 255, 255)

    image = Image.open(image_path).convert("RGB")
    langs = None # Replace with your languages or pass None (recommended to use None)
    recognition_predictor = RecognitionPredictor()
    detection_predictor = DetectionPredictor()

    surya_predictions = recognition_predictor([image], [langs], detection_predictor)[0].text_lines

    if len(surya_predictions) > 0:
        draw = ImageDraw.Draw(image)
        for line in surya_predictions:
                bbox = line.bbox
                original_text = line.text
                bbox = [
                    bbox[0] - 1,
                    bbox[1] - 1,
                    bbox[2] + 1,
                    bbox[3] + 1
                ] # padding

                # Перевод текста
                translated_text = ta.translate(
                                    source_lang=source_lang,
                                    target_lang=target_lang,
                                    source_text=original_text,
                                    do_reflection=do_reflection,
                                    country=country,
                                )

                draw.rectangle(bbox, fill=white_color)
                text_position = (
                    bbox[0] ,
                    (bbox[1] + bbox[3]) // 2
                )
                draw.text(text_position, translated_text, fill="black", font=font)

        image.save(image_path)
        print(f"Изображение {os.name(image_path)} успешно переведено")
    else:
        print("На изображении не был обнаружен текст.")


def markdown_to_word(markdown_content, word_file):
    # Converting Markdown to HTML
    html_content = markdown.markdown(markdown_content)

    # Creating a new Word Document
    doc = Document()

    # Converting HTML to text and add it to the Word Document
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Adding content to the Word Document
    for element in soup:
        if element.name == 'h1':
            doc.add_heading(element.text, level=1)
        elif element.name == 'h2':
            doc.add_heading(element.text, level=2)
        elif element.name == 'h3':
            doc.add_heading(element.text, level=3)
        elif element.name == 'p':
            paragraph = doc.add_paragraph()
            for child in element.children:
                if child.name == 'strong':
                    paragraph.add_run(child.text).bold = True
                elif child.name == 'em':
                    paragraph.add_run(child.text).italic = True
                else:
                    paragraph.add_run(child)
        elif element.name == 'ul':
            for li in element.find_all('li'):
                doc.add_paragraph(li.text, style='List Bullet')
        elif element.name == 'ol':
            for li in element.find_all('li'):
                doc.add_paragraph(li.text, style='List Number')
    
    doc.save(word_file)

def replace_pagination(text):
    # Регулярное выражение для поиска паттерна пагинации {число}----
    pagination_pattern = r'\{\d+\}-+'
    # Убираем остатки промта
    text = text.replace("</TRANSLATE_THIS_TO_RUSSIAN>", "")

    # Находим все вхождения пагинации
    matches = re.findall(pagination_pattern, text)

    # Если есть хотя бы одно вхождение пагинации
    if matches:
        # Удаляем первое вхождение пагинации
        first_pagination = matches[0]
        text = text.replace(first_pagination, "", 1)

    # Заменяем найденные паттерны на HTML-тег для разрыва страницы
    replaced_text = re.sub(pagination_pattern, '\n<div style="page-break-after: always"></div>\n', text)
    
    return replaced_text


def measure_time(func):
    """Декоратор для измерения времени выполнения функции."""
    def wrapper(*args, **kwargs):
        start_time = time.time()
        result = func(*args, **kwargs)
        end_time = time.time()
        
        elapsed_time = end_time - start_time
        minutes = int(elapsed_time // 60)
        seconds = int(elapsed_time % 60)
        formatted_time = f"{minutes:02d}:{seconds:02d}"
        
        print(f"Функция '{func.__name__}' выполнилась за {formatted_time}")
        return result
    return wrapper


def parse_arguments():
    parser = argparse.ArgumentParser(description="Перевод документов")

    ##### Аргументы для перевода #####
    parser.add_argument(
        "-i", "--input_language",
        required=False,
        help="Язык ввода (например, 'en', 'fr', 'de')"
    )
    parser.add_argument(
        "-o", "--output_language",
        required=True,
        help="Язык вывода (например, 'ru')"
    )
    parser.add_argument(
        "--do_reflection",
        action="store_true",
        help="Многостадийный перевод для повышения качества, True\False"
    )
    parser.add_argument(
        "-c", "--country",
        required=False,
        help="Страна (например, 'USA', 'France', 'Russia'). Необходимо для уточнения многостадийного перевода"
    )
    parser.add_argument(
        "-f", "--input_file",
        required=True,
        help="Путь к входному файлу"
    )
    parser.add_argument(
        "-d", "--output_dir",
        required=True,
        help="Путь к выходной директории"
    )

    ##### Аргументы для OCR #####
    parser.add_argument(
        "--paginate",
        action="store_true",
        help="Разбивать на страницы."
    )
    parser.add_argument(
        "--force_ocr",
        action="store_true",
        help="Форсировать OCR (может улучшить результаты OCR)."
    )
    parser.add_argument(
        "--page_range",
        required=False,
        help="Номера страниц документа для обработки, e.g. 1-5, 7"
    )

    # Парсим аргументы
    args = parser.parse_args()

    return args

def save_as_docx(translated_text, output_path):
    markdown_to_word(translated_text, output_path)

def save_as_pdf(translated_text, output_path):
    try:
        md2pdf(pdf_file_path = output_path,
            md_content=translated_text,
            base_url=os.path.dirname(output_path),
            css_file_path="md2pdf.css"
        )
    except Exception as e:
        print(f"Ошибка при обработке pdf: {e}")

def save_as_pptx(translated_text, output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.placeholders[0].text = translated_text
    prs.save(output_path)

def save_as_xlsx(translated_text, output_path):
    wb = Workbook()
    ws = wb.active
    ws.append([translated_text])
    wb.save(output_path)

def save_as_html(translated_text, output_path):
    soup = BeautifulSoup(f"<html><body><p>{translated_text}</p></body></html>", "html.parser")
    with open(output_path, "w", encoding="utf-8") as file:
        file.write(str(soup))

def save_as_epub(translated_text, output_path):
    book = epub.EpubBook()
    book.set_identifier('id123456')
    book.set_title('Translated Document')
    book.set_language('ru')
    c1 = epub.EpubHtml(title='Chapter 1', file_name='chap_01.xhtml', lang='ru')
    c1.content = f"<h1>Translated Document</h1><p>{translated_text}</p>"
    book.add_item(c1)
    book.spine = ['nav', c1]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(output_path, book)